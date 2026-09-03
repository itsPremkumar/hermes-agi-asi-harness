# -*- coding: utf-8 -*-
"""Agent Search Lite — Document Intelligence Layer.

Extracts content from PDF, DOCX, PPTX, XLSX files.

Copyright (c) 2026 Agent Search Lite Contributors.
MIT License. See LICENSE for details.
"""

from __future__ import annotations

import io
import logging
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional

import httpx

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# PDF Extraction
# ---------------------------------------------------------------------------

def extract_pdf(url_or_path: str, max_pages: int = 50) -> Dict[str, Any]:
    """Extract text and metadata from PDF.
    
    Uses multiple fallback methods:
    1. PyMuPDF (fitz) - best quality
    2. pdfplumber - table extraction
    3. pymupdf4llm - LLM-optimized markdown
    """
    result = {
        "url": url_or_path,
        "title": "",
        "text": "",
        "pages": 0,
        "metadata": {},
        "tables": [],
    }
    
    # Download if URL
    pdf_path = url_or_path
    temp_file = None
    
    if url_or_path.startswith(("http://", "https://")):
        try:
            resp = httpx.get(url_or_path, timeout=60, follow_redirects=True)
            resp.raise_for_status()
            temp_file = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
            temp_file.write(resp.content)
            temp_file.close()
            pdf_path = temp_file.name
        except Exception as exc:
            result["error"] = f"Download failed: {exc}"
            return result
    
    # Try PyMuPDF first
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(pdf_path)
        result["pages"] = len(doc)
        result["metadata"] = dict(doc.metadata) if doc.metadata else {}
        result["title"] = result["metadata"].get("title", "")
        
        text_parts = []
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            text_parts.append(page.get_text())
        
        result["text"] = "\n\n".join(text_parts)
        doc.close()
        
    except ImportError:
        # Fallback to pdfplumber
        try:
            import pdfplumber
            
            with pdfplumber.open(pdf_path) as pdf:
                result["pages"] = len(pdf.pages)
                result["metadata"] = dict(pdf.metadata) if pdf.metadata else {}
                
                text_parts = []
                for i, page in enumerate(pdf.pages):
                    if i >= max_pages:
                        break
                    text = page.extract_text() or ""
                    text_parts.append(text)
                    
                    # Extract tables
                    tables = page.extract_tables()
                    for table in tables:
                        result["tables"].append(table)
                
                result["text"] = "\n\n".join(text_parts)
                
        except ImportError:
            result["error"] = "No PDF library available. Install PyMuPDF or pdfplumber."
    
    except Exception as exc:
        result["error"] = f"PDF extraction failed: {exc}"
    
    finally:
        # Cleanup temp file
        if temp_file and os.path.exists(temp_file.name):
            os.unlink(temp_file.name)
    
    return result


# ---------------------------------------------------------------------------
# DOCX Extraction
# ---------------------------------------------------------------------------

def extract_docx(url_or_path: str) -> Dict[str, Any]:
    """Extract text and metadata from DOCX files."""
    result = {
        "url": url_or_path,
        "title": "",
        "text": "",
        "paragraphs": [],
        "tables": [],
        "metadata": {},
    }
    
    docx_path = url_or_path
    temp_file = None
    
    if url_or_path.startswith(("http://", "https://")):
        try:
            resp = httpx.get(url_or_path, timeout=60, follow_redirects=True)
            resp.raise_for_status()
            temp_file = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
            temp_file.write(resp.content)
            temp_file.close()
            docx_path = temp_file.name
        except Exception as exc:
            result["error"] = f"Download failed: {exc}"
            return result
    
    try:
        from docx import Document
        
        doc = Document(docx_path)
        
        # Extract metadata
        result["metadata"] = {
            "author": doc.core_properties.author,
            "title": doc.core_properties.title,
            "subject": doc.core_properties.subject,
            "keywords": doc.core_properties.keywords,
            "created": str(doc.core_properties.created) if doc.core_properties.created else "",
            "modified": str(doc.core_properties.modified) if doc.core_properties.modified else "",
        }
        result["title"] = result["metadata"].get("title", "")
        
        # Extract paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                result["paragraphs"].append(text)
        
        result["text"] = "\n\n".join(result["paragraphs"])
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            result["tables"].append(table_data)
        
    except ImportError:
        result["error"] = "python-docx not installed"
    except Exception as exc:
        result["error"] = f"DOCX extraction failed: {exc}"
    
    finally:
        if temp_file and os.path.exists(temp_file.name):
            os.unlink(temp_file.name)
    
    return result


# ---------------------------------------------------------------------------
# PPTX Extraction
# ---------------------------------------------------------------------------

def extract_pptx(url_or_path: str) -> Dict[str, Any]:
    """Extract text and metadata from PPTX files."""
    result = {
        "url": url_or_path,
        "title": "",
        "text": "",
        "slides": [],
        "metadata": {},
    }
    
    pptx_path = url_or_path
    temp_file = None
    
    if url_or_path.startswith(("http://", "https://")):
        try:
            resp = httpx.get(url_or_path, timeout=60, follow_redirects=True)
            resp.raise_for_status()
            temp_file = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
            temp_file.write(resp.content)
            temp_file.close()
            pptx_path = temp_file.name
        except Exception as exc:
            result["error"] = f"Download failed: {exc}"
            return result
    
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        
        # Extract metadata
        result["metadata"] = {
            "author": prs.core_properties.author,
            "title": prs.core_properties.title,
            "subject": prs.core_properties.subject,
            "slide_count": len(prs.slides),
        }
        result["title"] = result["metadata"].get("title", "")
        
        # Extract text from slides
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                result["slides"].append({
                    "slide_number": i + 1,
                    "text": "\n".join(slide_text),
                })
        
        result["text"] = "\n\n".join(s["text"] for s in result["slides"])
        
    except ImportError:
        result["error"] = "python-pptx not installed"
    except Exception as exc:
        result["error"] = f"PPTX extraction failed: {exc}"
    
    finally:
        if temp_file and os.path.exists(temp_file.name):
            os.unlink(temp_file.name)
    
    return result


# ---------------------------------------------------------------------------
# XLSX Extraction
# ---------------------------------------------------------------------------

def extract_xlsx(url_or_path: str) -> Dict[str, Any]:
    """Extract text and metadata from XLSX files."""
    result = {
        "url": url_or_path,
        "title": "",
        "text": "",
        "sheets": [],
        "metadata": {},
    }
    
    xlsx_path = url_or_path
    temp_file = None
    
    if url_or_path.startswith(("http://", "https://")):
        try:
            resp = httpx.get(url_or_path, timeout=60, follow_redirects=True)
            resp.raise_for_status()
            temp_file = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
            temp_file.write(resp.content)
            temp_file.close()
            xlsx_path = temp_file.name
        except Exception as exc:
            result["error"] = f"Download failed: {exc}"
            return result
    
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
        
        result["metadata"] = {
            "sheet_names": wb.sheetnames,
            "sheet_count": len(wb.sheetnames),
        }
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(row_text):
                    rows.append(row_text)
            
            if rows:
                result["sheets"].append({
                    "name": sheet_name,
                    "rows": rows[:100],  # Limit rows
                })
        
        # Create text representation
        text_parts = []
        for sheet in result["sheets"]:
            text_parts.append(f"## {sheet['name']}")
            for row in sheet["rows"][:20]:
                text_parts.append(" | ".join(str(c) for c in row))
        
        result["text"] = "\n".join(text_parts)
        
    except ImportError:
        result["error"] = "openpyxl not installed"
    except Exception as exc:
        result["error"] = f"XLSX extraction failed: {exc}"
    
    finally:
        if temp_file and os.path.exists(temp_file.name):
            os.unlink(temp_file.name)
    
    return result


# ---------------------------------------------------------------------------
# Universal Document Extractor
# ---------------------------------------------------------------------------

def extract_document(url_or_path: str) -> Dict[str, Any]:
    """Universal document extractor that auto-detects file type."""
    # Detect file type from extension or content-type
    lower_url = url_or_path.lower()
    
    if lower_url.endswith(".pdf"):
        return extract_pdf(url_or_path)
    elif lower_url.endswith(".docx"):
        return extract_docx(url_or_path)
    elif lower_url.endswith(".pptx"):
        return extract_pptx(url_or_path)
    elif lower_url.endswith(".xlsx") or lower_url.endswith(".xls"):
        return extract_xlsx(url_or_path)
    else:
        # Try to detect from content-type
        try:
            if url_or_path.startswith(("http://", "https://")):
                resp = httpx.head(url_or_path, timeout=15, follow_redirects=True)
                content_type = resp.headers.get("content-type", "")
                
                if "pdf" in content_type:
                    return extract_pdf(url_or_path)
                elif "word" in content_type or "docx" in content_type:
                    return extract_docx(url_or_path)
                elif "presentation" in content_type or "pptx" in content_type:
                    return extract_pptx(url_or_path)
                elif "sheet" in content_type or "xlsx" in content_type:
                    return extract_xlsx(url_or_path)
            
            return {"url": url_or_path, "error": "Unknown document type"}
            
        except Exception:
            return {"url": url_or_path, "error": "Unknown document type"}
